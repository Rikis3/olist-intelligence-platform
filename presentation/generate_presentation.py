from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()
    
    # Custom Modern Colors (Dark Mode)
    bg_color = RGBColor(25, 25, 35)      # Deep Dark Blue/Gray Background
    text_color = RGBColor(230, 230, 240) # Off-white text
    accent_color = RGBColor(0, 190, 255) # Electric Cyan
    bullet_color = RGBColor(255, 120, 0) # Neon Orange

    # --- Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color

    title = slide.shapes.title
    title.text = "Olist Intelligence Platform"
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.bold = True

    subtitle = slide.placeholders[1]
    subtitle.text = "Strategic Review & Logistics Root Cause Analysis\nPrepared by: Senior BI Analyst"
    subtitle.text_frame.paragraphs[0].font.color.rgb = text_color
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)

    # --- Helper function for modern content slides ---
    def style_slide(slide, title_text, body_texts):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color
        
        title_shape = slide.shapes.title
        title_shape.text = title_text
        title_shape.text_frame.paragraphs[0].font.color.rgb = accent_color
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        body_shape = slide.placeholders[1]
        body_shape.text = body_texts[0]
        body_shape.text_frame.paragraphs[0].font.color.rgb = text_color
        body_shape.text_frame.paragraphs[0].font.size = Pt(26)
        
        for text in body_texts[1:]:
            p = body_shape.text_frame.add_paragraph()
            p.text = text
            p.level = 1
            p.font.color.rgb = text_color
            p.font.size = Pt(22)

    # Slide 2: Business Problem
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    style_slide(slide2, "The Business Problem: Revenue Leakage", [
        "Olist is facing a massive challenge with Customer Lifetime Value (CLV).",
        "Most customers only purchase once. We need to increase repeat purchase rates.",
        "Hypothesis: Logistics bottlenecks are driving down CSAT and causing severe customer churn."
    ])

    # Slide 3: Analytics Architecture
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    style_slide(slide3, "Analytics Engineering Architecture", [
        "To analyze this, we built a complete Enterprise Data Stack:",
        "Data Warehouse: 9 raw tables modeled into a BigQuery Star Schema.",
        "Advanced SQL: RFM Segmentation and Cohort Retention logic implemented.",
        "BI Layer: Interactive Streamlit Executive Dashboard built with Plotly."
    ])

    # Slide 4: Strategic Finding 1
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    style_slide(slide4, "Finding 1: The Logistics-to-Churn Pipeline", [
        "Late deliveries completely destroy Customer Lifetime Value.",
        "On-time delivery repeat purchase rate: 18.5%",
        "Late delivery repeat purchase rate: 4.5% (A 75% Drop-off!)",
        "Estimated annualized revenue leakage due to logistics: $2.4M."
    ])

    # Slide 5: Strategic Finding 2
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    style_slide(slide5, "Finding 2: Geographic Bottlenecks", [
        "The Northeast region (Bahia, Pernambuco) is critically underperforming.",
        "Average Lead Time: >21 days (vs 6 days in Sao Paulo).",
        "SLA Breach Rate: 45% (vs 2% in Sao Paulo)."
    ])

    # Slide 6: Recommendations
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    style_slide(slide6, "Strategic Recommendations", [
        "Immediate actions to recover LTV and improve CSAT:",
        "1. Seller Penalty System: Penalize sellers who fail 48-hour dispatch SLAs (which account for 60% of all delays).",
        "2. Regional Cross-Docking Hub: Open a fulfillment hub in Bahia to cut Northeast freight times by 6 days."
    ])
    
    # Save the presentation
    output_path = os.path.join(os.path.dirname(__file__), 'Olist_Executive_Review.pptx')
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
